"""Belge ayrıştırıcı — PDF, DOCX, PPTX, TXT, MD dosyalarından düz metin çıkarır."""
import io
from pathlib import Path


def parse_document(content: bytes, filename: str) -> str:
    """Dosya içeriğini ve adını alır, düz metin döndürür."""
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(content)
    elif ext == ".docx":
        return _parse_docx(content)
    elif ext == ".pptx":
        return _parse_pptx(content)
    elif ext in (".txt", ".md"):
        return content.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Desteklenmeyen dosya formatı: {ext}")


def _parse_pdf(content: bytes) -> str:
    try:
        import PyPDF2

        reader = PyPDF2.PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(p for p in pages if p.strip())
    except Exception as exc:
        raise ValueError(f"PDF ayrıştırma hatası: {exc}") from exc


def _parse_docx(content: bytes) -> str:
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:
        raise ValueError(f"DOCX ayrıştırma hatası: {exc}") from exc


def _parse_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides_text: list[str] = []
        for slide in prs.slides:
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if slide_lines:
                slides_text.append("\n".join(slide_lines))
        return "\n\n".join(slides_text)
    except Exception as exc:
        raise ValueError(f"PPTX ayrıştırma hatası: {exc}") from exc
